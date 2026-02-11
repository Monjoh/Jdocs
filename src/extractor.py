"""File metadata extraction engine for jDocs.

Extracts text content and metadata from supported file types:
- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)
- Images (.png, .jpg, .jpeg)
- CSV (.csv) — first 100 rows, column names, row count
- Code files (.py, .js, .java, .ts, .c, .cpp, .html, .css, .json, .xml, .md, .txt)
"""

import csv
import io
from pathlib import Path
from zipfile import BadZipFile

from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from PIL.ExifTags import TAGS
from pptx import Presentation


# Maximum file size for text extraction (10 MB)
MAX_TEXT_SIZE = 10 * 1024 * 1024

# File extensions recognized as plain-text code/config files
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
    ".html", ".css", ".json", ".xml", ".yaml", ".yml",
    ".md", ".txt", ".sh", ".bat", ".ps1",
    ".rb", ".go", ".rs", ".swift", ".kt",
}


def extract(file_path: str | Path) -> dict:
    """Extract text and metadata from a file.

    Returns a dict with keys:
        - file_name: original filename
        - file_type: extension (e.g. ".docx")
        - size_bytes: file size
        - text: extracted text content
        - metadata: dict of type-specific metadata
        - error: error message if extraction failed, else None
    """
    path = Path(file_path)

    if not path.exists():
        return _error_result(path, "File not found")

    if not path.is_file():
        return _error_result(path, "Not a file (may be a directory)")

    ext = path.suffix.lower()
    base = {
        "file_name": path.name,
        "file_type": ext,
        "size_bytes": path.stat().st_size,
        "text": "",
        "metadata": {},
        "error": None,
    }

    try:
        if ext == ".docx":
            _extract_docx(path, base)
        elif ext == ".xlsx":
            _extract_xlsx(path, base)
        elif ext == ".pptx":
            _extract_pptx(path, base)
        elif ext == ".csv":
            _extract_csv(path, base)
        elif ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}:
            _extract_image(path, base)
        elif ext in CODE_EXTENSIONS:
            _extract_code(path, base)
        else:
            base["metadata"]["note"] = "Unsupported file type — no extraction performed"
    except BadZipFile:
        base["error"] = "This file appears to be password-protected or corrupted."
    except Exception as e:
        msg = str(e).lower()
        if "package not found" in msg or "not a zip file" in msg:
            base["error"] = "This file appears to be password-protected or corrupted."
        else:
            base["error"] = f"Extraction failed: {e}"

    return base


def _error_result(path: Path, message: str) -> dict:
    return {
        "file_name": path.name,
        "file_type": path.suffix.lower(),
        "size_bytes": 0,
        "text": "",
        "metadata": {},
        "error": message,
    }


def _extract_docx(path: Path, result: dict):
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    result["text"] = "\n".join(paragraphs)

    props = doc.core_properties
    result["metadata"] = {
        "author": props.author or "",
        "title": props.title or "",
        "paragraph_count": len(paragraphs),
    }


def _extract_xlsx(path: Path, result: dict):
    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheet_info = []
    all_text = []
    max_preview_rows = 100

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        row_count = 0
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            row_count += 1
            # Only collect text from the first N rows for preview
            if i < max_preview_rows:
                for cell in row:
                    if cell is not None:
                        all_text.append(str(cell))
        sheet_info.append({"name": sheet_name, "row_count": row_count})

    wb.close()
    result["text"] = "\n".join(all_text)
    result["metadata"] = {
        "sheet_count": len(wb.sheetnames),
        "sheets": sheet_info,
    }


def _extract_pptx(path: Path, result: dict):
    prs = Presentation(str(path))
    slide_texts = []

    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        slide_texts.append("\n".join(parts))

    result["text"] = "\n\n".join(slide_texts)

    props = prs.core_properties
    result["metadata"] = {
        "author": props.author or "",
        "title": props.title or "",
        "slide_count": len(prs.slides),
    }


def _extract_image(path: Path, result: dict):
    img = Image.open(path)
    result["metadata"] = {
        "width": img.width,
        "height": img.height,
        "format": img.format,
        "mode": img.mode,
    }

    # Try to extract EXIF data
    exif_data = {}
    try:
        raw_exif = img.getexif()
        if raw_exif:
            for tag_id, value in raw_exif.items():
                tag_name = TAGS.get(tag_id, str(tag_id))
                # Only keep simple serializable values
                if isinstance(value, (str, int, float)):
                    exif_data[tag_name] = value
    except Exception:
        pass

    if exif_data:
        result["metadata"]["exif"] = exif_data

    img.close()


def _extract_csv(path: Path, result: dict):
    """Extract column names and first 100 rows from a CSV file."""
    max_preview_rows = 100
    size = path.stat().st_size
    truncated = size > MAX_TEXT_SIZE
    if truncated:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read(MAX_TEXT_SIZE)
    else:
        raw = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(raw))

    rows = []
    for i, row in enumerate(reader):
        if i > max_preview_rows:
            break
        rows.append(row)

    # Total row count by counting newlines (fast, avoids reading all rows via csv)
    total_rows = raw.count("\n")

    columns = rows[0] if rows else []
    sample_text = []
    for row in rows[:max_preview_rows]:
        sample_text.append(", ".join(row))

    result["text"] = "\n".join(sample_text)
    result["metadata"] = {
        "columns": columns,
        "column_count": len(columns),
        "total_rows": total_rows,
        "preview_rows": min(max_preview_rows, len(rows)),
    }


def _extract_code(path: Path, result: dict):
    size = path.stat().st_size
    truncated = size > MAX_TEXT_SIZE
    if truncated:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read(MAX_TEXT_SIZE)
    else:
        text = path.read_text(encoding="utf-8", errors="replace")
    result["text"] = text
    result["metadata"] = {
        "line_count": text.count("\n") + 1 if text else 0,
        "char_count": len(text),
    }
    if truncated:
        result["metadata"]["truncated"] = True
        result["metadata"]["note"] = f"File truncated to first 10 MB (full size: {size:,} bytes)"
